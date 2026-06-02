"""Production de livrables Word (.docx) et PowerPoint (.pptx).

Le modèle passe son contenu en Markdown — on le convertit côté Python. Plus
simple et plus fiable que de demander au modèle de produire des objets
structurés multiples.

Word :
    # Titre 1
    ## Titre 2
    ### Titre 3
    Paragraphe normal.
    - puce
    1. liste numérotée

PowerPoint (slides séparées par une ligne `---`) :
    # Titre de la slide
    - puce 1
    - puce 2
    > Note du présentateur (n'apparaît pas sur la slide)
"""

from __future__ import annotations

import re
from datetime import date
from pathlib import Path

from anthropic import beta_tool

try:
    from docx import Document
except ImportError:  # pragma: no cover
    Document = None  # type: ignore

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:  # pragma: no cover
    Presentation = None  # type: ignore

from ..config import load_config
from .documents import VALID_FOLDERS, _safe_filename, _root


def _resolve_path(folder: str, filename: str, ext: str) -> Path:
    base, _, _ = filename.rpartition(".")
    base = base or filename
    base = re.sub(r"[^A-Za-z0-9._-]+", "-", base).strip("-") or "document"
    if not base.startswith(date.today().isoformat()):
        base = f"{date.today().isoformat()}_{base}"
    return _root() / folder / f"{base}.{ext}"


def _render_docx(path: Path, contenu_markdown: str, titre: str | None) -> None:
    doc = Document()
    if titre:
        doc.add_heading(titre, level=0)
    for raw in contenu_markdown.splitlines():
        line = raw.rstrip()
        if not line:
            doc.add_paragraph("")
            continue
        m = re.match(r"^(#{1,3})\s+(.+)$", line)
        if m:
            doc.add_heading(m.group(2), level=len(m.group(1)))
            continue
        if line.startswith("- ") or line.startswith("* "):
            doc.add_paragraph(line[2:], style="List Bullet")
            continue
        if re.match(r"^\d+\.\s", line):
            doc.add_paragraph(re.sub(r"^\d+\.\s", "", line), style="List Number")
            continue
        doc.add_paragraph(line)
    doc.save(path)


def _render_pptx(path: Path, slides_markdown: str, titre: str | None) -> int:
    prs = Presentation()
    # Slide de titre
    if titre:
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = titre
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = f"{date.today().isoformat()}"

    blocks = re.split(r"^\s*---\s*$", slides_markdown, flags=re.MULTILINE)
    count = 0
    for block in blocks:
        lines = [l for l in block.splitlines() if l.strip()]
        if not lines:
            continue
        # 1ʳᵉ ligne non vide = titre
        title_line = re.sub(r"^#+\s*", "", lines[0]).strip()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title_line
        bullets: list[str] = []
        notes: list[str] = []
        for line in lines[1:]:
            stripped = line.strip()
            if stripped.startswith("> "):
                notes.append(stripped[2:])
                continue
            if stripped.startswith("- ") or stripped.startswith("* "):
                bullets.append(stripped[2:])
            else:
                bullets.append(stripped)
        body_ph = slide.placeholders[1].text_frame
        body_ph.text = bullets[0] if bullets else ""
        for b in bullets[1:]:
            p = body_ph.add_paragraph()
            p.text = b
            p.font.size = Pt(18)
        if notes:
            slide.notes_slide.notes_text_frame.text = "\n".join(notes)
        count += 1

    prs.save(path)
    return count


@beta_tool
def save_word_document(
    folder: str,
    filename: str,
    titre: str,
    contenu_markdown: str,
) -> str:
    """Sauvegarde un document Word (.docx) dans SM_PROJECT/<folder>/.

    Le contenu est passé en Markdown — il est converti en Word automatiquement :
    `# Titre`, `## Sous-titre`, `### Sous-sous-titre`, paragraphes simples,
    listes `- puce` ou `1. élément`.

    Utilise-le pour les livrables structurés destinés à un client :
    propositions commerciales, comptes rendus, supports écrits, AO.

    Args:
        folder: section cible (01_Strategie, 02_Prospection, 03_Missions, 04_Clients,
            05_Formation_MS_Project, 06_LinkedIn, 07_Administratif, 08_Finances,
            09_Portfolio, 10_Opportunites).
        filename: nom court du fichier (sans extension, .docx ajouté).
        titre: titre principal qui apparaît en tête du document.
        contenu_markdown: contenu complet en Markdown.
    """
    if Document is None:
        return "python-docx non installé. Lance: pip install python-docx"
    if folder not in VALID_FOLDERS:
        return f"Folder invalide. Valeurs : {', '.join(VALID_FOLDERS)}."
    path = _resolve_path(folder, filename, "docx")
    _render_docx(path, contenu_markdown, titre)
    rel = path.relative_to(load_config().data_dir)
    return f"Word créé : {rel} ({path.stat().st_size} octets)."


@beta_tool
def save_powerpoint(
    folder: str,
    filename: str,
    titre: str,
    slides_markdown: str,
) -> str:
    """Sauvegarde une présentation PowerPoint (.pptx) dans SM_PROJECT/<folder>/.

    Chaque slide est séparée par une ligne `---`. Dans chaque slide :
    - 1ʳᵉ ligne = titre (préfixe `#` toléré)
    - lignes `- texte` ou `* texte` = puces
    - lignes commençant par `> ` = notes du présentateur

    Utilise-le pour les pitchs, propositions visuelles, plans de formation,
    supports de RDV de découverte.

    Args:
        folder: section cible (voir save_word_document pour les valeurs).
        filename: nom court (sans extension, .pptx ajouté).
        titre: titre de la slide de couverture.
        slides_markdown: ensemble des slides, séparées par `---`.

    Exemple de slides_markdown :
        # Pourquoi un PMO externalisé
        - Vision projet partagée CODIR
        - Reporting fiabilisé
        > Insister sur le gain de visibilité direction
        ---
        # Notre offre
        - 3 jours/mois
        - Engagement révisable
    """
    if Presentation is None:
        return "python-pptx non installé. Lance: pip install python-pptx"
    if folder not in VALID_FOLDERS:
        return f"Folder invalide. Valeurs : {', '.join(VALID_FOLDERS)}."
    path = _resolve_path(folder, filename, "pptx")
    n = _render_pptx(path, slides_markdown, titre)
    rel = path.relative_to(load_config().data_dir)
    return f"PowerPoint créé : {rel} ({n} slide(s), {path.stat().st_size} octets)."


TOOLS = [save_word_document, save_powerpoint]
